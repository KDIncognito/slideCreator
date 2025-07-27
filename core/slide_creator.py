import os
import sys
import tempfile
import shutil
from pathlib import Path

# Add the project root to Python path for imports
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

# Project imports (using absolute imports)
from documentHandling.convert_pdf_to_image import PDFToImageConverter
from utils.logger import get_logger

# Try to import optional components (they may not exist yet)
try:
    from interfaces.base_interfaces import ISlideCreator, ProcessingResult
    from validation.input_validator import InputValidator
    from integration.content_visual_bridge import ContentVisualBridge
    from presentation.powerpoint_generator import PowerPointGenerator
    FULL_SYSTEM_AVAILABLE = True
except ImportError as e:
    print(f"Optional components not available: {e}")
    FULL_SYSTEM_AVAILABLE = False
    
    # Create minimal compatibility classes
    class ISlideCreator:
        pass
    
    class ProcessingResult:
        def __init__(self, success, errors=None, warnings=None, metadata=None):
            self.success = success
            self.errors = errors or []
            self.warnings = warnings or []
            self.metadata = metadata or {}


class SlideCreator(ISlideCreator):
    """Main orchestrator for the PDF to PowerPoint conversion process."""
    
    def __init__(self):
        self.temp_dir = None
        self.logger = get_logger()
        
        # Initialize components if available
        if FULL_SYSTEM_AVAILABLE:
            try:
                self.validator = InputValidator()
                self.bridge = ContentVisualBridge()
                self.ppt_generator = PowerPointGenerator()
            except Exception as e:
                print(f"Error initializing full system components: {e}")
                self.validator = None
                self.bridge = None
                self.ppt_generator = None
        else:
            self.validator = None
            self.bridge = None
            self.ppt_generator = None
    
    def validate_environment(self) -> ProcessingResult:
        """
        Validate that all required environment variables and dependencies are available.
        
        Returns:
            ProcessingResult with validation outcome
        """
        errors = []
        warnings = []
        
        # Check OpenAI API key
        if not os.getenv('gpt_key'):
            errors.append("OpenAI API key not found. Set 'gpt_key' environment variable.")
        
        # Check required modules
        required_modules = [
            ('openai', 'OpenAI API client'),
            ('fitz', 'PyMuPDF for PDF processing'),
            ('PIL', 'Pillow for image processing'),
            ('pptx', 'python-pptx for PowerPoint generation')
        ]
        
        for module_name, description in required_modules:
            try:
                __import__(module_name)
            except ImportError:
                errors.append(f"Missing required module: {description}")
        
        # Check system resources
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                test_file = Path(temp_dir) / "test.txt"
                test_file.write_text("test")
        except Exception as e:
            warnings.append(f"Temporary directory access issue: {str(e)}")
        
        # Check if full system is available
        if not FULL_SYSTEM_AVAILABLE:
            warnings.append("Some system components are not available - running in basic mode")
        
        return ProcessingResult(
            success=len(errors) == 0,
            errors=errors,
            warnings=warnings,
            metadata={'validation_type': 'environment', 'full_system': FULL_SYSTEM_AVAILABLE}
        )
    
    def get_conversion_info(self, pdf_path: str) -> dict:
        """
        Get information about a PDF file without converting it.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Dictionary containing PDF information
        """
        pdf_file = Path(pdf_path)
        
        if not pdf_file.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        if not pdf_file.suffix.lower() == '.pdf':
            raise ValueError("Input file must be a PDF")
        
        return {
            'file_name': pdf_file.name,
            'file_size_mb': pdf_file.stat().st_size / (1024 * 1024),
            'file_path': str(pdf_file.absolute()),
            'expected_output': str(pdf_file.with_suffix('.pptx')),
            'is_readable': pdf_file.exists() and os.access(pdf_file, os.R_OK)
        }
    
    def convert_pdf_to_slides(self, pdf_path: str, output_path: str = None) -> str:
        """
        Convert a PDF file to a PowerPoint presentation.
        
        Args:
            pdf_path: Path to the input PDF file
            output_path: Optional path for output PowerPoint file
            
        Returns:
            Path to the generated PowerPoint file
        """
        # Validate input
        pdf_file = Path(pdf_path)
        if not pdf_file.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        
        if not pdf_file.suffix.lower() == '.pdf':
            raise ValueError("Input file must be a PDF")
        
        # Set up output path
        if output_path is None:
            output_path = str(pdf_file.with_suffix('.pptx'))
        
        self.logger.info(f"(>>>) Converting {pdf_file.name} to PowerPoint presentation...")
        
        try:
            # Step 1: Convert PDF to images
            self.logger.info("(###) Converting PDF to images...")
            pdf_data = self._process_pdf(pdf_path)
            
            # Step 2: Create PowerPoint presentation
            if FULL_SYSTEM_AVAILABLE and self.bridge:
                self.logger.info("(^^^) Analyzing content with AI...")
                slide_structure = self._analyze_content_with_ai(pdf_data)
                final_path = self._generate_advanced_powerpoint(slide_structure, output_path)
            else:
                self.logger.info("(@@@) Creating basic PowerPoint presentation...")
                final_path = self._create_basic_powerpoint(pdf_data, output_path)
            
            self.logger.info(f"(!!!) Conversion complete! Presentation saved to: {final_path}")
            return final_path
            
        except Exception as e:
            self.logger.error(f"(XXX) Error during conversion: {str(e)}")
            raise
        finally:
            self._cleanup()
    
    def _process_pdf(self, pdf_path: str) -> dict:
        """
        Process PDF to extract images.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Dictionary containing image paths and basic info
        """
        # Create temporary directory for images
        self.temp_dir = tempfile.mkdtemp(prefix="slidecreator_")
        self.logger.debug(f"Created temporary directory: {self.temp_dir}")
        
        # Convert PDF to images
        converter = PDFToImageConverter()
        image_paths = converter.convert_pdf_to_images(pdf_path, self.temp_dir)
        
        self.logger.debug(f"Processed {len(image_paths)} pages from PDF")
        
        return {
            'image_paths': image_paths,
            'total_pages': len(image_paths),
            'text_by_page': {page_num: f"[Page {page_num} content]" for page_num in image_paths.keys()}
        }
    
    def _analyze_content_with_ai(self, pdf_data: dict) -> dict:
        """
        Analyze content using the full AI pipeline.
        
        Args:
            pdf_data: Dictionary containing extracted data
            
        Returns:
            Dictionary containing slide structure and analysis results
        """
        if not self.bridge:
            raise RuntimeError("AI analysis not available - ContentVisualBridge not initialized")
        
        self.logger.debug("Starting content analysis with ContentVisualBridge")
        
        result = self.bridge.process_complete_document(
            extracted_text_by_page=pdf_data['text_by_page'],
            image_paths_by_page=pdf_data['image_paths']
        )
        
        # Validate that we got slides
        if 'slides' not in result or 'slides' not in result['slides']:
            raise RuntimeError("Failed to generate slide structure from PDF content")
        
        slides_count = len(result['slides']['slides'])
        self.logger.debug(f"Generated {slides_count} slides with AI")
        
        return result
    
    def _generate_advanced_powerpoint(self, slide_structure: dict, output_path: str) -> str:
        """
        Generate PowerPoint using the advanced generator.
        
        Args:
            slide_structure: Dictionary containing slide data and structure
            output_path: Path where to save the PowerPoint file
            
        Returns:
            Path to the generated PowerPoint file
        """
        if not self.ppt_generator:
            raise RuntimeError("Advanced PowerPoint generator not available")
        
        slides_data = slide_structure['slides']
        
        self.logger.debug(f"Creating advanced PowerPoint file: {output_path}")
        
        final_path = self.ppt_generator.create_presentation(
            slides_data=slides_data,
            output_path=output_path,
            visual_mappings=slide_structure.get('visual_mappings', [])
        )
        
        return final_path
    
    def _create_basic_powerpoint(self, pdf_data: dict, output_path: str) -> str:
        """
        Create a basic PowerPoint presentation.
        This fallback method works without the full AI pipeline.
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")
        
        self.logger.info("Creating basic PowerPoint presentation...")
        
        # Create new presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Generated Presentation"
        subtitle.text = f"Converted from PDF\n{pdf_data['total_pages']} pages processed"
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Content slides - one per PDF page
        for page_num in range(1, pdf_data['total_pages'] + 1):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = f"Page {page_num} Content"
            
            # Style title
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Add content
            tf = body_shape.text_frame
            tf.text = f"Content extracted from PDF page {page_num}"
            
            p = tf.add_paragraph()
            p.text = "(AI analysis and detailed content extraction available with full system)"
            p.level = 1
            
            # Style body text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(24)
        
        # Save presentation
        prs.save(output_path)
        
        self.logger.info(f"Basic PowerPoint created with {pdf_data['total_pages']} content slides")
        
        return output_path
    
    def _cleanup(self):
        """Clean up temporary files and directories."""
        if self.temp_dir and os.path.exists(self.temp_dir):
            self.logger.debug(f"Cleaning up temporary directory: {self.temp_dir}")
            try:
                shutil.rmtree(self.temp_dir)
            except Exception as e:
                self.logger.warning(f"(!!!) Failed to cleanup temporary directory: {str(e)}")
