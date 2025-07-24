from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json
import glob

class WritePpt:
    def __init__(self):
        self.presentation = Presentation()
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        self.image_directory = "images_and_files/"
        
    def find_image_file(self, image_id):
        """Find image file with given ID in various formats"""
        if not image_id or image_id == "NA":
            return None
            
        # Check for different image formats
        formats = ['png', 'jpg', 'jpeg', 'gif', 'bmp']
        
        for fmt in formats:
            image_path = os.path.join(self.image_directory, f"{image_id}.{fmt}")
            if os.path.exists(image_path):
                return image_path
                
        # Also try wildcard search in case of different naming
        pattern = os.path.join(self.image_directory, f"{image_id}.*")
        matches = glob.glob(pattern)
        if matches:
            return matches[0]
            
        return None
    
    def parse_recommended_placement(self, recommended_placement):
        """Parse the specific recommended_placement enum from your LLM prompts"""
        if not recommended_placement or recommended_placement == "NA":
            return {"left": Inches(6), "top": Inches(2), "width": Inches(3)}
        
        placement = recommended_placement.upper()
        
        # Based on your prompt's enum values:
        # 'FULL_SLIDE', 'LEFT_HALF', 'RIGHT_HALF', 'TOP_HALF', 'BOTTOM_HALF', 'BACKGROUND_OVERLAY'
        
        if placement == 'FULL_SLIDE':
            return {
                "left": Inches(0.5),
                "top": Inches(1.5),
                "width": Inches(9),
                "height": Inches(5)
            }
        elif placement == 'LEFT_HALF':
            return {
                "left": Inches(0.5),
                "top": Inches(2),
                "width": Inches(4.5),
                "height": None
            }
        elif placement == 'RIGHT_HALF':
            return {
                "left": Inches(5.5),
                "top": Inches(2),
                "width": Inches(4),
                "height": None
            }
        elif placement == 'TOP_HALF':
            return {
                "left": Inches(2),
                "top": Inches(1.5),
                "width": Inches(6),
                "height": Inches(2.5)
            }
        elif placement == 'BOTTOM_HALF':
            return {
                "left": Inches(2),
                "top": Inches(4.5),
                "width": Inches(6),
                "height": Inches(2.5)
            }
        elif placement == 'BACKGROUND_OVERLAY':
            return {
                "left": Inches(1),
                "top": Inches(1),
                "width": Inches(8),
                "height": Inches(6),
                "send_to_back": True
            }
        else:
            # Default to right half
            return {
                "left": Inches(5.5),
                "top": Inches(2),
                "width": Inches(4),
                "height": None
            }
    
    def get_layout_by_type(self, suggested_layout_type):
        """Get PowerPoint layout based on your LLM's suggested_layout_type enum"""
        layout_map = {
            'TITLE_ONLY': 0,  # Title slide
            'TITLE_BULLETS': 1,  # Title and Content
            'TITLE_CONTENT_BULLETS': 1,  # Title and Content
            'TITLE_IMAGE_RIGHT': 1,  # Title and Content (we'll add image)
            'TITLE_IMAGE_LEFT': 1,  # Title and Content (we'll add image)
            'TITLE_IMAGE_FULL': 5,  # Blank slide (we'll add title and image)
            'TITLE_TWO_COLUMNS': 3,  # Two Column Content
            'SECTION_HEADER': 2,  # Section Header
            'TITLE_BULLETS_IMAGE_RIGHT': 1,  # Title and Content (we'll add image)
        }
        
        return layout_map.get(suggested_layout_type, 1)  # Default to Title and Content
    
    def create_title_slide(self, title, subtitle=None):
        """Create a title slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            
        return slide
    
    def create_content_slide_from_llm_structure(self, slide_data):
        """Create a slide based on your exact LLM output structure"""
        print(f"DEBUG: Creating slide with structure: {slide_data.keys()}")
        
        # Extract data from your LLM structure
        title = slide_data.get('title', 'Untitled Slide')
        content = slide_data.get('content', [])
        suggested_layout_type = slide_data.get('suggested_layout_type', 'TITLE_BULLETS')
        generated_visual_placeholder = slide_data.get('generated_visual_placeholder')
        
        # Get appropriate layout
        layout_index = self.get_layout_by_type(suggested_layout_type)
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Handle image placement based on your LLM's structure
        image_added = False
        content_area_adjusted = False
        
        if generated_visual_placeholder and generated_visual_placeholder.get('need_image', False):
            image_id = generated_visual_placeholder.get('image_placeholder_id')
            recommended_placement = generated_visual_placeholder.get('recommended_placement', 'RIGHT_HALF')
            caption = generated_visual_placeholder.get('caption', '')
            description = generated_visual_placeholder.get('description_for_audience', '')
            
            print(f"DEBUG: Processing image - ID: {image_id}, Placement: {recommended_placement}")
            
            if image_id and image_id != "NA":
                image_path = self.find_image_file(image_id)
                
                if image_path:
                    try:
                        pos_params = self.parse_recommended_placement(recommended_placement)
                        
                        # Add image with your LLM's specific positioning
                        if pos_params.get('height'):
                            image_shape = slide.shapes.add_picture(
                                image_path, 
                                pos_params['left'], 
                                pos_params['top'], 
                                width=pos_params['width'],
                                height=pos_params['height']
                            )
                        else:
                            image_shape = slide.shapes.add_picture(
                                image_path, 
                                pos_params['left'], 
                                pos_params['top'], 
                                width=pos_params['width']
                            )
                        
                        # Handle background overlay
                        if pos_params.get('send_to_back'):
                            # Move to back (PowerPoint doesn't have direct API, but we can note it)
                            pass
                        
                        image_added = True
                        print(f"✓ Image added: {image_path} at placement: {recommended_placement}")
                        
                        # Adjust content area based on image placement
                        if recommended_placement in ['LEFT_HALF', 'RIGHT_HALF']:
                            content_area_adjusted = True
                        
                    except Exception as e:
                        print(f"XX Error adding image {image_path}: {e}")
                else:
                    print(f"⚠️ Image not found: {image_id}")
        
        # Add content with layout adjustment
        if content and len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            # Adjust content positioning if image affects layout
            if content_area_adjusted and generated_visual_placeholder:
                placement = generated_visual_placeholder.get('recommended_placement', '')
                if placement == 'LEFT_HALF':
                    # Content goes to right side
                    content_shape.left = Inches(5.5)
                    content_shape.width = Inches(4)
                elif placement == 'RIGHT_HALF':
                    # Content stays on left but gets narrower
                    content_shape.width = Inches(5)
            
            # Add bullet points
            for i, point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        return slide
    
    def create_slides_from_content(self, slide_content_dict):
        """Create slides from your handler's exact slide content dictionary structure"""
        print(f"DEBUG: Creating slides from content type: {type(slide_content_dict)}")
        
        if not isinstance(slide_content_dict, dict):
            print(f"Error: slide_content must be a dictionary, got {type(slide_content_dict)}")
            return False
            
        slides_data = slide_content_dict.get('slides', [])
        if not slides_data:
            print("Error: No slides found in content")
            print(f"Available keys: {list(slide_content_dict.keys())}")
            return False
        
        print(f"DEBUG: Found {len(slides_data)} slides to create")
        
        # Create title slide
        main_title = slide_content_dict.get('title', 'Generated Presentation')
        subtitle = slide_content_dict.get('subtitle', 'Created from PDF Content')
        self.create_title_slide(main_title, subtitle)
        
        # Create content slides using your exact LLM structure
        for i, slide_data in enumerate(slides_data):
            print(f"DEBUG: Processing slide {i+1} with keys: {list(slide_data.keys())}")
            self.create_content_slide_from_llm_structure(slide_data)
            
        return True
    
    def save_presentation(self, filename="generated_presentation.pptx"):
        """Save the presentation to a file"""
        try:
            if not filename.endswith('.pptx'):
                filename += '.pptx'
                
            self.presentation.save(filename)
            print(f"✓ Presentation saved as: {filename}")
            return True
        except Exception as e:
            print(f"XX Error saving presentation: {e}")
            return False
    
    def get_slide_count(self):
        """Get the number of slides in the presentation"""
        return len(self.presentation.slides)

# Convenience function for easy use
def create_powerpoint_from_content(slide_content, output_filename="generated_slides.pptx"):
    """Create a PowerPoint presentation from slide content using your exact LLM format"""
    print(f"DEBUG: Starting PowerPoint creation with content type: {type(slide_content)}")
    
    ppt_writer = WritePpt()
    
    if ppt_writer.create_slides_from_content(slide_content):
        success = ppt_writer.save_presentation(output_filename)
        if success:
            print(f"✓ PowerPoint created successfully: {output_filename}")
            print(f"✓ Total slides: {ppt_writer.get_slide_count()}")
            return True
    
    print("XX Failed to create PowerPoint presentation")
    return False

# if __name__ == "__main__":
#     # Test with your exact LLM structure
#     test_content = {
#         "title": "Test Presentation",
#         "slides": [
#             {
#                 "title": "Quantum-AI Synergy",
#                 "content": [
#                     "Quantum computing enhances AI capabilities",
#                     "Exponential speedup for complex algorithms",
#                     "Revolutionary applications in optimization"
#                 ],
#                 "generated_visual_placeholder": {
#                     "image_placeholder_id": "slide_01_visual_01",
#                     "concept_id_link": "concept_001",
#                     "description_for_audience": "Figure 1: Conceptual Diagram of Quantum-AI Synergy",
#                     "recommended_placement": "RIGHT_HALF",
#                     "need_image": True,
#                     "caption": "Visualizing the fusion of quantum mechanics and machine learning.",
#                     "source_text_for_visual": "Quantum machine learning explores the fusion of quantum computing with machine learning to enhance data processing and pattern recognition beyond classical limits."
#                 },
#                 "suggested_layout_type": "TITLE_BULLETS_IMAGE_RIGHT"
#             }
#         ]
#     }
    
#     create_powerpoint_from_content(test_content, "test_llm_format.pptx")